#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def build_ppt(input_dir: Path, output_file: Path) -> int:
    images = sorted(input_dir.glob("*.png"))
    if not images:
        raise SystemExit(f"No PNG files found in: {input_dir}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for image in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(image), 0, 0, width=slide_w, height=slide_h)

    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    return len(images)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Build a PPTX from PNG screenshots (one image per slide)."
    )
    parser.add_argument(
        "--input-dir",
        type=Path,
        default=Path("screenshots/openclaw-complete-scene-ends-4k"),
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("out/OpenClawComplete-SceneEnds-4k.pptx"),
    )
    args = parser.parse_args()

    count = build_ppt(args.input_dir, args.output)
    print(f"Created {args.output} with {count} slides.")


if __name__ == "__main__":
    main()
